"""
PowerPoint CRUD — Créer, lire, modifier et supprimer des fichiers PPTX.

Catégorie: file
Mode: sync

Actions:
    create  → Génère un PPTX depuis des données structurées (slides, texte, images)
    read    → Extrait le contenu d'un PPTX (slides, texte, notes)
    update  → Modifie un PPTX (ajouter des slides, remplacer du texte)
    delete  → Supprime un fichier PPTX du stockage
"""

from __future__ import annotations

import io
from typing import Any

from app.framework.base import BaseTool
from app.framework.schemas import (
    ToolErrorCode,
    ToolExample,
    ToolExecutionMode,
    ToolMetadata,
    ToolParameter,
    ToolResult,
)


class PowerpointCrud(BaseTool):
    """CRUD complet pour les fichiers PPTX via MinIO."""

    @property
    def metadata(self) -> ToolMetadata:
        return ToolMetadata(
            slug="powerpoint-crud",
            name="PowerPoint CRUD",
            description="Créer, lire, modifier et supprimer des fichiers PowerPoint (PPTX)",
            version="1.0.0",
            category="file",
            execution_mode=ToolExecutionMode.SYNC,
            timeout_seconds=30,
            input_schema=[
                ToolParameter(name="action", type="string", required=True,
                              description="Action: create, read, update, delete"),
                ToolParameter(name="storage_key", type="string",
                              description="Chemin MinIO du fichier"),
                ToolParameter(name="data", type="object",
                              description="Contenu: {slides: [{title, content, notes, layout}]}"),
                ToolParameter(name="options", type="object",
                              description="Options: {width, height}"),
            ],
            output_schema=[
                ToolParameter(name="storage_key", type="string"),
                ToolParameter(name="slides", type="array",
                              description="Liste des slides avec contenu"),
                ToolParameter(name="slide_count", type="integer"),
                ToolParameter(name="text", type="string",
                              description="Texte complet extrait"),
            ],
            examples=[
                ToolExample(
                    description="Créer une présentation",
                    input={
                        "action": "create",
                        "storage_key": "presentations/demo.pptx",
                        "data": {
                            "slides": [
                                {"title": "Bienvenue", "content": "Présentation NOVA2", "layout": "title"},
                                {"title": "Agenda", "content": "1. Intro\n2. Démo\n3. Q&A", "layout": "content"},
                                {"title": "Merci!", "content": "Questions?", "layout": "title"},
                            ],
                        },
                    },
                    output={"storage_key": "presentations/demo.pptx", "slide_count": 3},
                ),
            ],
            tags=["file", "powerpoint", "pptx", "crud", "office", "presentation"],
        )

    async def execute(self, params: dict[str, Any], context) -> ToolResult:
        action = params.get("action", "")
        if action == "create":
            return await self._create(params, context)
        elif action == "read":
            return await self._read(params, context)
        elif action == "update":
            return await self._update(params, context)
        elif action == "delete":
            return await self._delete(params, context)
        else:
            return self.error(
                f"Action inconnue: '{action}'. Actions: create, read, update, delete",
                ToolErrorCode.INVALID_PARAMS,
            )

    async def _create(self, params: dict[str, Any], context) -> ToolResult:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        storage_key = params.get("storage_key", "")
        data = params.get("data", {})

        if not storage_key:
            return self.error("storage_key requis pour create", ToolErrorCode.INVALID_PARAMS)

        prs = Presentation()

        slides_data = data.get("slides", [])
        for slide_data in slides_data:
            title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            notes = slide_data.get("notes", "")
            layout_type = slide_data.get("layout", "content")

            # Choisir le layout
            if layout_type == "title":
                layout = prs.slide_layouts[0]  # Title Slide
            elif layout_type == "section":
                layout = prs.slide_layouts[2]  # Section Header
            elif layout_type == "blank":
                layout = prs.slide_layouts[6]  # Blank
            else:
                layout = prs.slide_layouts[1]  # Title and Content

            slide = prs.slides.add_slide(layout)

            # Titre
            if slide.shapes.title and title:
                slide.shapes.title.text = title

            # Contenu
            if content and len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = content

            # Notes
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        await context.storage.put(
            storage_key, buffer.getvalue(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        return self.success({
            "storage_key": storage_key,
            "slide_count": len(slides_data),
        })

    async def _read(self, params: dict[str, Any], context) -> ToolResult:
        from pptx import Presentation

        storage_key = params.get("storage_key", "")
        if not storage_key:
            return self.error("storage_key requis pour read", ToolErrorCode.INVALID_PARAMS)

        file_data = await context.storage.get(storage_key)
        if file_data is None:
            return self.error(f"Fichier introuvable: {storage_key}", ToolErrorCode.FILE_NOT_FOUND)

        try:
            prs = Presentation(io.BytesIO(file_data))
        except Exception as e:
            return self.error(f"Fichier PPTX invalide: {e}", ToolErrorCode.PROCESSING_ERROR)

        slides = []
        full_text_parts = []

        for idx, slide in enumerate(prs.slides):
            slide_info = {"index": idx, "shapes": []}

            # Titre
            if slide.shapes.title:
                slide_info["title"] = slide.shapes.title.text
                full_text_parts.append(slide.shapes.title.text)

            # Contenu textuel de toutes les shapes
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if text.strip():
                        texts.append(text)
                        full_text_parts.append(text)
                        slide_info["shapes"].append({
                            "type": "text",
                            "text": text,
                            "name": shape.name,
                        })
                elif shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    slide_info["shapes"].append({
                        "type": "table",
                        "data": table_data,
                        "name": shape.name,
                    })

            slide_info["content"] = "\n".join(texts)

            # Notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                slide_info["notes"] = notes_text

            slides.append(slide_info)

        return self.success({
            "storage_key": storage_key,
            "slides": slides,
            "slide_count": len(slides),
            "text": "\n\n".join(full_text_parts),
        })

    async def _update(self, params: dict[str, Any], context) -> ToolResult:
        from pptx import Presentation

        storage_key = params.get("storage_key", "")
        data = params.get("data", {})

        if not storage_key:
            return self.error("storage_key requis pour update", ToolErrorCode.INVALID_PARAMS)

        file_data = await context.storage.get(storage_key)
        if file_data is None:
            return self.error(f"Fichier introuvable: {storage_key}", ToolErrorCode.FILE_NOT_FOUND)

        try:
            prs = Presentation(io.BytesIO(file_data))
        except Exception as e:
            return self.error(f"Fichier PPTX invalide: {e}", ToolErrorCode.PROCESSING_ERROR)

        # Remplacer du texte dans toutes les slides
        replacements = data.get("replace", {})
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for old_text, new_text in replacements.items():
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)

        # Ajouter des slides
        add_slides = data.get("add_slides", [])
        for slide_data in add_slides:
            title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            layout_type = slide_data.get("layout", "content")

            if layout_type == "title":
                layout = prs.slide_layouts[0]
            elif layout_type == "blank":
                layout = prs.slide_layouts[6]
            else:
                layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(layout)
            if slide.shapes.title and title:
                slide.shapes.title.text = title
            if content and len(slide.placeholders) > 1:
                slide.placeholders[1].text_frame.text = content

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        await context.storage.put(
            storage_key, buffer.getvalue(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        return self.success({
            "storage_key": storage_key,
            "slide_count": len(prs.slides),
        })

    async def _delete(self, params: dict[str, Any], context) -> ToolResult:
        storage_key = params.get("storage_key", "")
        if not storage_key:
            return self.error("storage_key requis pour delete", ToolErrorCode.INVALID_PARAMS)

        deleted = await context.storage.delete(storage_key)
        if not deleted:
            return self.error(f"Fichier introuvable: {storage_key}", ToolErrorCode.FILE_NOT_FOUND)

        return self.success({"storage_key": storage_key, "deleted": True})
